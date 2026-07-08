#!/usr/bin/env python3
"""生成「朝夕」生活秩序管理助手 - 开题答辩PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 颜色方案
PRIMARY = RGBColor(0x25, 0x63, 0xEB)  # 蓝色主色
DARK = RGBColor(0x1E, 0x29, 0x3B)      # 深色文字
GRAY = RGBColor(0x64, 0x74, 0x8B)       # 灰色
LIGHT_BG = RGBColor(0xF1, 0xF5, 0xF9)   # 浅背景
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 辅助函数 ──

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT, font_name='PingFang SC'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf

def add_shape(slide, left, top, width, height, color, shape_type=MSO_SHAPE.RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_rounded_card(slide, left, top, width, height, color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_bullet_frame(slide, left, top, width, height, items, font_size=16, color=DARK):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'PingFang SC'
        p.space_after = Pt(8)
        p.level = 0
    return tf

# ========== 第1页：封面 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, PRIMARY)

# 装饰条
add_shape(slide, 0, 6.8, 13.333, 0.7, RGBColor(0x1D, 0x4E, 0xD8))

add_textbox(slide, 1.5, 1.8, 10, 1.2, '开源项目答辩', font_size=22, color=RGBColor(0xBF, 0xDB, 0xFE), bold=False)
add_textbox(slide, 1.5, 2.4, 10, 1.5, '朝夕 — 生活秩序管理助手', font_size=48, color=WHITE, bold=True)
add_textbox(slide, 1.5, 3.6, 10, 0.8, '基于 HarmonyOS + Dify AI 的个人生活管理应用', font_size=20, color=RGBColor(0xDB, 0xEA, 0xFE))

# 底部信息
add_textbox(slide, 1.5, 5.5, 5, 0.5, '答辩人：XXX', font_size=16, color=WHITE)
add_textbox(slide, 1.5, 5.9, 5, 0.5, '指导老师：XXX', font_size=16, color=WHITE)
add_textbox(slide, 1.5, 6.3, 5, 0.5, '2026年6月', font_size=14, color=RGBColor(0x93, 0xC5, 0xFD))

# ========== 第2页：目录 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, 0, 0, 13.333, 0.08, PRIMARY)

add_textbox(slide, 1.0, 0.5, 5, 0.8, '目  录', font_size=36, color=DARK, bold=True)
add_textbox(slide, 1.0, 1.2, 8, 0.5, 'CONTENTS', font_size=14, color=GRAY)

toc_items = [
    ('01', '项目背景与意义'),
    ('02', '项目目标与用户群体'),
    ('03', '系统总体架构'),
    ('04', '核心功能模块'),
    ('05', 'Dify AI 工作流设计'),
    ('06', 'UI 界面展示'),
    ('07', '技术栈与开发环境'),
    ('08', '项目计划与进度'),
]

for i, (num, title) in enumerate(toc_items):
    y = 1.6 + i * 0.75
    add_shape(slide, 1.0, y + 0.05, 0.06, 0.35, PRIMARY)
    add_textbox(slide, 1.3, y, 0.6, 0.4, num, font_size=24, color=PRIMARY, bold=True)
    add_textbox(slide, 2.0, y + 0.02, 8, 0.4, title, font_size=17, color=DARK)

# ========== 第3页：项目背景 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, 0, 0, 13.333, 0.08, PRIMARY)
add_textbox(slide, 1.0, 0.5, 8, 0.8, '01  项目背景与意义', font_size=32, color=DARK, bold=True)

# 左栏 - 背景
add_rounded_card(slide, 0.8, 1.8, 5.5, 2.5)
add_textbox(slide, 1.1, 2.0, 4, 0.5, '社会背景', font_size=20, color=PRIMARY, bold=True)
bg_items = [
    '现代大学生普遍存在作息不规律、拖延等问题',
    '缺乏有效的自我管理工具，传统待办App功能单一',
    'AI技术发展为个性化生活指导提供了可能',
    '鸿蒙生态崛起，校园应用场景丰富',
]
add_bullet_frame(slide, 1.1, 2.6, 5.0, 1.5, bg_items, font_size=14, color=DARK)

# 右栏 - 意义
add_rounded_card(slide, 7.0, 1.8, 5.5, 2.5)
add_textbox(slide, 7.3, 2.0, 4, 0.5, '项目意义', font_size=20, color=PRIMARY, bold=True)
meaning_items = [
    '将AI大模型能力落地到日常生活管理场景',
    '通过游戏化成就系统提升用户自律动力',
    '探索 Dify 低代码AI工作流的应用边界',
    '为鸿蒙生态贡献一款实用型应用',
]
add_bullet_frame(slide, 7.3, 2.6, 5.0, 1.5, meaning_items, font_size=14, color=DARK)

# 底部 - 痛点分析
add_textbox(slide, 0.8, 4.8, 12, 0.6, '用户痛点分析', font_size=20, color=PRIMARY, bold=True)

pain_points = [
    ('😴', '作息混乱', '入睡时间不固定\n缺乏规律作息监督'),
    ('📋', '任务拖延', '待办事项堆积\n缺乏优先级管理'),
    ('📝', '缺乏反思', '没有记录习惯\n无法回溯成长'),
    ('🧠', '缺乏指导', '需要个性化\nAI建议和反馈'),
]

for i, (icon, title, desc) in enumerate(pain_points):
    x = 0.8 + i * 3.1
    add_rounded_card(slide, x, 5.5, 2.8, 1.6)
    add_textbox(slide, x + 0.3, 5.6, 2, 0.4, f'{icon} {title}', font_size=16, color=DARK, bold=True)
    add_textbox(slide, x + 0.3, 6.0, 2.2, 0.9, desc, font_size=12, color=GRAY)

# ========== 第4页：项目目标 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, 0, 0, 13.333, 0.08, PRIMARY)
add_textbox(slide, 1.0, 0.5, 8, 0.8, '02  项目目标与用户群体', font_size=32, color=DARK, bold=True)

# 目标
add_textbox(slide, 1.0, 1.8, 12, 0.6, '项目目标', font_size=22, color=PRIMARY, bold=True)

goals = [
    ('核心目标', '构建一款集成 AI 智能分析的生活秩序管理应用'),
    ('功能目标', '实现作息打卡、待办管理、日记记录、成就系统四大核心模块'),
    ('技术目标', '完成 HarmonyOS 前端与 Dify AI 后端的联调，实现智能复盘与建议'),
    ('体验目标', '通过游戏化设计提升用户粘性，打造简洁美观的交互界面'),
]
for i, (label, desc) in enumerate(goals):
    y = 2.5 + i * 0.7
    add_shape(slide, 1.0, y + 0.08, 1.6, 0.45, PRIMARY)
    add_textbox(slide, 1.1, y + 0.12, 1.5, 0.4, label, font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, 2.9, y + 0.1, 9, 0.4, desc, font_size=15, color=DARK)

# 用户群体
add_textbox(slide, 1.0, 5.3, 12, 0.6, '目标用户群体', font_size=22, color=PRIMARY, bold=True)

users = [
    ('👨‍🎓', '大学生', '希望改善作息\n提升学习效率'),
    ('👩‍💻', '职场新人', '需要任务管理\n平衡工作生活'),
    ('📱', '鸿蒙用户', '追求简洁体验\n注重隐私安全'),
]
for i, (icon, title, desc) in enumerate(users):
    x = 1.0 + i * 3.8
    add_rounded_card(slide, x, 5.9, 3.3, 1.3)
    add_textbox(slide, x + 0.3, 6.0, 3, 0.4, f'{icon}  {title}', font_size=18, color=DARK, bold=True)
    add_textbox(slide, x + 0.3, 6.4, 3, 0.6, desc, font_size=13, color=GRAY)

# ========== 第5页：系统架构 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, 0, 0, 13.333, 0.08, PRIMARY)
add_textbox(slide, 1.0, 0.5, 8, 0.8, '03  系统总体架构', font_size=32, color=DARK, bold=True)

# 三层架构
layers = [
    ('表现层  HarmonyOS App', PRIMARY, ['首页仪表盘', '待办清单', '日记日历', '个人中心', '成就徽章']),
    ('服务层  Dify AI 平台', RGBColor(0x3B, 0x82, 0xF6), ['数据预处理', '知识库检索', '联网搜索', 'LLM对话生成', 'JSON格式化输出']),
    ('数据层  本地 + 云端', RGBColor(0x63, 0x66, 0xF1), ['SQLite本地存储', 'iCloud/云存储同步', '用户行为日志', '知识库向量化']),
]

for i, (title, color, items) in enumerate(layers):
    y = 1.8 + i * 1.8
    # 层标签
    add_shape(slide, 1.0, y, 2.8, 1.4, color)
    add_textbox(slide, 1.2, y + 0.45, 2.5, 0.6, title, font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # 箭头连接
    if i < 2:
        add_textbox(slide, 2.4, y + 1.35, 0.5, 0.5, '⬇', font_size=20, color=GRAY, align=PP_ALIGN.CENTER)
    # 内容卡片
    for j, item in enumerate(items):
        x = 4.3 + j * 1.75
        add_rounded_card(slide, x, y + 0.15, 1.55, 1.1, LIGHT_BG)
        add_textbox(slide, x + 0.1, y + 0.35, 1.4, 0.6, item, font_size=12, color=DARK, align=PP_ALIGN.CENTER)

# 底部说明
add_textbox(slide, 1.0, 7.0, 11, 0.4, '▸ App通过RESTful API与Dify平台通信，Dify工作流处理AI请求并返回结构化JSON数据', font_size=13, color=GRAY)

# ========== 第6页：核心功能 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, 0, 0, 13.333, 0.08, PRIMARY)
add_textbox(slide, 1.0, 0.5, 8, 0.8, '04  核心功能模块', font_size=32, color=DARK, bold=True)

modules = [
    ('🌙', '作息打卡', '记录每日入睡/起床时间\n自动计算睡眠时长\nAI分析睡眠质量', PRIMARY),
    ('✅', '待办管理', '创建/编辑/删除待办\n优先级分类管理\nAI智能拆解大任务', ACCENT_GREEN),
    ('📝', '日记记录', '日历视图浏览\n心情标记（4种状态）\nAI生成每日小结', ACCENT_ORANGE),
    ('🏆', '成就系统', '等级成长体系\n12种成就徽章\n连续打卡统计', RGBColor(0x8B, 0x5C, 0xF6)),
]

for i, (icon, title, desc, color) in enumerate(modules):
    x = 0.6 + i * 3.15
    add_rounded_card(slide, x, 1.8, 2.9, 3.8)
    # 顶部色条
    add_shape(slide, x, 1.8, 2.9, 0.08, color)
    add_textbox(slide, x + 0.3, 2.1, 2, 0.5, f'{icon}  {title}', font_size=22, color=color, bold=True)
    add_textbox(slide, x + 0.3, 2.7, 2.3, 2.0, desc, font_size=14, color=DARK)

# 底部亮点
add_rounded_card(slide, 0.6, 6.0, 12.1, 1.1, LIGHT_BG)
add_textbox(slide, 1.0, 6.15, 11, 0.8, '💡 亮点：游戏化成就系统驱动用户行为改变 — 完成打卡/待办/日记获得成就点，解锁徽章，提升等级', font_size=15, color=DARK)

# ========== 第7页：Dify AI工作流 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, 0, 0, 13.333, 0.08, PRIMARY)
add_textbox(slide, 1.0, 0.5, 8, 0.8, '05  Dify AI 工作流设计', font_size=32, color=DARK, bold=True)

add_textbox(slide, 1.0, 1.5, 11, 0.5, '工作流管道（5个节点）', font_size=18, color=GRAY)

nodes = [
    ('① 数据预处理', '清洗用户输入\n标准化格式\n提取关键信息', PRIMARY),
    ('② 知识库检索', '匹配历史数据\n检索作息模式\n相似场景对比', RGBColor(0x3B, 0x82, 0xF6)),
    ('③ 联网搜索', '补充外部知识\n获取最新资讯\n扩展建议范围', ACCENT_GREEN),
    ('④ LLM生成', '生成个性化建议\n自然语言总结\n多轮对话能力', ACCENT_ORANGE),
    ('⑤ JSON格式化', '结构化输出\n字段标准化\nAPI友好格式', RGBColor(0x8B, 0x5C, 0xF6)),
]

for i, (title, desc, color) in enumerate(nodes):
    x = 0.5 + i * 2.55
    add_rounded_card(slide, x, 2.1, 2.3, 3.0)
    add_shape(slide, x, 2.1, 2.3, 0.06, color)
    add_textbox(slide, x + 0.2, 2.3, 1.9, 0.5, title, font_size=16, color=color, bold=True)
    add_textbox(slide, x + 0.2, 2.9, 1.9, 1.8, desc, font_size=13, color=DARK)
    # 箭头
    if i < 4:
        add_textbox(slide, x + 2.3, 3.3, 0.3, 0.4, '→', font_size=22, color=GRAY, align=PP_ALIGN.CENTER)

# 应用场景
add_textbox(slide, 1.0, 5.5, 12, 0.5, 'AI 应用场景', font_size=20, color=PRIMARY, bold=True)

scenes = [
    'AI作息诊断：分析用户睡眠模式，给出改善建议',
    'AI每日小结：自动总结当日完成情况，生成日记草稿',
    'AI任务拆解：将复杂任务拆分为可执行的子任务',
    'AI月度复盘：基于历史数据生成月度趋势分析报告',
]

for i, scene in enumerate(scenes):
    y = 6.1 + i * 0.35
    add_textbox(slide, 1.2, y, 11, 0.35, f'▸ {scene}', font_size=14, color=DARK)

# ========== 第8页：UI展示 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_shape(slide, 0, 0, 13.333, 0.08, PRIMARY)
add_textbox(slide, 1.0, 0.5, 8, 0.8, '06  UI 界面展示', font_size=32, color=DARK, bold=True)

ui_pages = [
    ('🏠', '首页', '问候栏 + 今日成就\n进度条 + 作息卡片\n待办概览 + AI功能区'),
    ('✅', '待办页', '日期切换 + 任务列表\n优先级标签 + 完成率\nAI拆解按钮'),
    ('📝', '日记页', '月度日历 + 心情标记\n日记列表 + 快速浏览\n悬浮添加按钮'),
    ('👤', '我的', '等级卡 + 进度条\n成就徽章墙\n功能入口列表'),
]

for i, (icon, title, desc) in enumerate(ui_pages):
    x = 0.5 + i * 3.2
    add_rounded_card(slide, x, 1.8, 2.9, 4.5, WHITE)
    add_textbox(slide, x + 0.3, 2.0, 2, 0.5, f'{icon}  {title}', font_size=20, color=PRIMARY, bold=True)
    add_textbox(slide, x + 0.3, 2.6, 2.3, 2.0, desc, font_size=13, color=DARK)
    # 手机框示意
    add_shape(slide, x + 0.5, 4.5, 1.9, 1.5, LIGHT_BG, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, x + 0.7, 5.0, 1.5, 0.5, '📱', font_size=30, color=GRAY, align=PP_ALIGN.CENTER)

add_textbox(slide, 0.8, 6.7, 12, 0.6, '▸ 采用HarmonyOS设计规范，Tailwind CSS风格配色，移动端优先设计', font_size=14, color=GRAY)

# ========== 第9页：技术栈 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, 0, 0, 13.333, 0.08, PRIMARY)
add_textbox(slide, 1.0, 0.5, 8, 0.8, '07  技术栈与开发环境', font_size=32, color=DARK, bold=True)

stacks = [
    ('前端技术', [
        'HarmonyOS ArkTS + ArkUI',
        'JavaScript/TypeScript',
        '响应式布局',
    ], PRIMARY),
    ('AI 平台', [
        'Dify 低代码AI平台',
        'OpenAI GPT-4 API',
        '向量知识库 (RAG)',
    ], RGBColor(0x3B, 0x82, 0xF6)),
    ('后端/数据', [
        'RESTful API',
        'SQLite / 云存储',
        'Node.js / Python',
    ], ACCENT_GREEN),
    ('开发工具', [
        'DevEco Studio',
        'Git 版本管理',
        'Postman API测试',
    ], RGBColor(0x8B, 0x5C, 0xF6)),
]

for i, (title, items, color) in enumerate(stacks):
    x = 0.5 + i * 3.2
    add_rounded_card(slide, x, 1.8, 2.9, 3.2)
    add_shape(slide, x, 1.8, 2.9, 0.06, color)
    add_textbox(slide, x + 0.3, 2.0, 2.3, 0.5, title, font_size=20, color=color, bold=True)
    add_bullet_frame(slide, x + 0.3, 2.6, 2.3, 2.0, items, font_size=14, color=DARK)

# 技术亮点
add_textbox(slide, 1.0, 5.4, 12, 0.5, '技术亮点', font_size=20, color=PRIMARY, bold=True)

highlights = [
    ('RAG知识库', '将用户历史数据向量化存储，实现个性化上下文检索'),
    ('低代码AI', 'Dify可视化编排工作流，降低AI集成门槛，快速迭代'),
    ('鸿蒙原生', '充分利用HarmonyOS分布式能力，适配多设备协同场景'),
]
for i, (title, desc) in enumerate(highlights):
    y = 5.9 + i * 0.45
    add_textbox(slide, 1.2, y, 2.5, 0.4, f'▸ {title}', font_size=14, color=PRIMARY, bold=True)
    add_textbox(slide, 3.8, y, 8, 0.4, desc, font_size=14, color=DARK)

# ========== 第10页：项目计划 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, 0, 0, 13.333, 0.08, PRIMARY)
add_textbox(slide, 1.0, 0.5, 8, 0.8, '08  项目计划与进度', font_size=32, color=DARK, bold=True)

# 甘特图风格
phases = [
    ('需求分析', '第1-2周', '✅ 已完成', ACCENT_GREEN, 1.0),
    ('UI原型设计', '第3-4周', '✅ 已完成', ACCENT_GREEN, 1.0),
    ('前端核心开发', '第5-8周', '🔄 进行中', PRIMARY, 0.6),
    ('Dify工作流', '第7-9周', '🔄 进行中', PRIMARY, 0.4),
    ('联调测试', '第10-11周', '⏳ 待开始', GRAY, 0),
    ('答辩与提交', '第12周', '⏳ 待开始', GRAY, 0),
]

for i, (phase, duration, status, color, progress) in enumerate(phases):
    y = 2.0 + i * 0.85
    add_textbox(slide, 1.0, y, 2.5, 0.4, phase, font_size=16, color=DARK, bold=True)
    add_textbox(slide, 3.6, y, 1.8, 0.4, duration, font_size=14, color=GRAY)
    # 进度条背景
    add_shape(slide, 5.6, y + 0.1, 5.0, 0.2, LIGHT_BG, MSO_SHAPE.ROUNDED_RECTANGLE)
    # 进度条填充
    if progress > 0:
        add_shape(slide, 5.6, y + 0.1, 5.0 * progress, 0.2, color, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, 10.9, y, 1.8, 0.4, status, font_size=14, color=color, bold=True)

# 风险与对策
add_textbox(slide, 1.0, 7.0, 12, 0.4, '风险：Dify API延迟影响体验 → 对策：本地缓存 + 异步加载  |  AI输出不可控 → 对策：工作流末尾JSON格式化约束输出', font_size=12, color=GRAY)

# ========== 第11页：致谢 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)
add_shape(slide, 0, 6.8, 13.333, 0.7, RGBColor(0x1D, 0x4E, 0xD8))

add_textbox(slide, 1.5, 2.5, 10, 1.5, '感谢聆听', font_size=52, color=WHITE, bold=True)
add_textbox(slide, 1.5, 3.8, 10, 0.8, '朝夕 — 生活秩序管理助手', font_size=24, color=RGBColor(0xBF, 0xDB, 0xFE))
add_textbox(slide, 1.5, 4.5, 10, 0.8, '欢迎提问与交流', font_size=20, color=WHITE)

add_textbox(slide, 1.5, 6.2, 5, 0.4, '联系方式：XXX', font_size=14, color=RGBColor(0x93, 0xC5, 0xFD))
add_textbox(slide, 1.5, 6.5, 5, 0.4, '项目仓库：github.com/xxx/zhaoxi', font_size=14, color=RGBColor(0x93, 0xC5, 0xFD))

# ── 保存 ──
output_path = '/Users/wenzhiyuan/Desktop/朝夕-开题答辩.pptx'
prs.save(output_path)
print(f'PPT已生成: {output_path}')

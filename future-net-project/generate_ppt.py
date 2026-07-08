#!/usr/bin/env python3
"""生成课程答辩 PPT — 白底、简洁、不溢出，4 页"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'results', '开题答辩.pptx')
os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 颜色 ──
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLACK   = RGBColor(0x1A, 0x1A, 0x1A)
BLUE    = RGBColor(0x1F, 0x6F, 0xEB)
GREEN   = RGBColor(0x18, 0x7D, 0x34)
RED     = RGBColor(0xD1, 0x24, 0x2F)
GRAY    = RGBColor(0x6B, 0x72, 0x7A)
LGRAY   = RGBColor(0xF3, 0xF4, 0xF6)
DGRAY   = RGBColor(0x3A, 0x3F, 0x47)

FONT   = 'Microsoft YaHei'
MONO   = 'Consolas'


def bg(s, c=WHITE):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = c


def txt(s, l, t, w, h, text, sz=12, c=BLACK, bold=False, align=PP_ALIGN.LEFT, font=FONT):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(sz); p.font.color.rgb = c
    p.font.bold = bold; p.font.name = font; p.alignment = align
    return tf


def title_bar(s, text, top=0.3):
    txt(s, 0.5, top, 8, 0.45, text, sz=26, c=BLUE, bold=True)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(top+0.52), Inches(2.2), Inches(0.035))
    ln.fill.solid(); ln.fill.fore_color.rgb = BLUE; ln.line.fill.background()


def card(s, l, t, w, h, title, lines, tc=BLUE, tsz=11, lsz=10):
    """灰色卡片"""
    bg_rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    bg_rect.fill.solid(); bg_rect.fill.fore_color.rgb = LGRAY
    bg_rect.line.color.rgb = RGBColor(0xDD, 0xDF, 0xE2); bg_rect.line.width = Pt(0.5)

    txt(s, l+0.12, t+0.06, w-0.24, 0.22, title, sz=tsz, c=tc, bold=True)

    body = s.shapes.add_textbox(Inches(l+0.12), Inches(t+0.32), Inches(w-0.24), Inches(h-0.4))
    bf = body.text_frame; bf.word_wrap = True
    bf.margin_left = Inches(0.02); bf.margin_right = Inches(0.02)
    bf.margin_top = Inches(0.02); bf.margin_bottom = Inches(0.02)
    for i, line in enumerate(lines):
        p = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        p.text = line; p.font.size = Pt(lsz); p.font.color.rgb = BLACK
        p.font.name = FONT; p.space_after = Pt(2)


def code(s, l, t, w, h, title, lines):
    """代码卡片"""
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    rect.fill.solid(); rect.fill.fore_color.rgb = RGBColor(0xFA, 0xFB, 0xFC)
    rect.line.color.rgb = RGBColor(0xDD, 0xDF, 0xE2); rect.line.width = Pt(0.5)

    txt(s, l+0.1, t+0.04, w-0.2, 0.2, title, sz=10, c=DGRAY, bold=True, font=MONO)

    body = s.shapes.add_textbox(Inches(l+0.1), Inches(t+0.28), Inches(w-0.2), Inches(h-0.35))
    bf = body.text_frame; bf.word_wrap = True
    bf.margin_left = Inches(0.02); bf.margin_right = Inches(0.02)
    bf.margin_top = Inches(0.02); bf.margin_bottom = Inches(0.02)
    for i, line in enumerate(lines):
        p = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        p.text = line; p.font.size = Pt(9); p.font.color.rgb = DGRAY
        p.font.name = MONO; p.space_after = Pt(1)


# ═══════════════════════════════════════════════════════
# Slide 1: 标题
# ═══════════════════════════════════════════════════════
s1 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s1)

txt(s1, 1.5, 1.8, 10, 1.2, '基于 SDN 与机器学习的\n智能流量分类系统',
    sz=42, c=BLUE, bold=True, align=PP_ALIGN.CENTER)
txt(s1, 1.5, 3.3, 10, 0.5, '未来网络技术课程实践',
    sz=20, c=GRAY, align=PP_ALIGN.CENTER)

ln = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.2), Inches(3.9), Inches(2.8), Inches(0.03))
ln.fill.solid(); ln.fill.fore_color.rgb = BLUE; ln.line.fill.background()

txt(s1, 1.5, 4.3, 10, 0.4, '姓名：XXX    学号：XXX', sz=18, c=BLACK, align=PP_ALIGN.CENTER)
txt(s1, 1.5, 4.8, 10, 0.4, '指导老师：朱轶', sz=16, c=GRAY, align=PP_ALIGN.CENTER)
txt(s1, 1.5, 6.3, 10, 0.3, '江苏大学    未来网络技术课程实践    2026',
    sz=12, c=GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════
# Slide 2: 作品简介
# ═══════════════════════════════════════════════════════
s2 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s2)
title_bar(s2, '作品简介')

# 概述（一行）
txt(s2, 0.5, 1.0, 12.3, 0.45,
    '用 Mininet 搭建 SDN 网络，Ryu 控制器通过 OpenFlow 采集 OVS 端口统计，调用独立的 AI 微服务实时识别流量类型（视频/网页/下载），并自动下发 QoS 策略。',
    sz=12, c=BLACK)

# 左列：系统架构 + 技术栈
card(s2, 0.5, 1.6, 4.1, 2.5, '系统架构与技术栈', [
    '数据层 · Mininet 2.3.0 + OVS 2.17',
    '  4主机+1交换机, MAC自学习L2转发(无路由), 子网10.0.0.0/24',
    '控制层 · Ryu 4.34 + OpenFlow 1.3',
    '  L2转发 · 每2s采集PortStats · 提取10维特征',
    'AI层 · Random Forest + Flask',
    '  200棵树, <1ms推理, 独立微服务(端口5001)',
    '流量 · Python socket UDP + curl + wget',
    '  三种协议特征 · 收发比区分 · 混合方案',
    '架构: OVS → Ryu → HTTP POST → Flask → predict()',
    '符合工业界 SDN两层分离 + AI微服务解耦模式',
], lsz=10)

# 中列：流量产生 + 特征
card(s2, 4.8, 1.6, 4.1, 2.5, '流量产生与 10 维特征', [
    'h1 · Python socket UDP → h4:9999',
    '  1200B/包, ~200pps, 无回包 → 收发比≈∞   → video',
    'h2 · curl 间歇 HTTP → h4:8000/',
    '  550B/包, ~30pps, TCP双向 → 收发比≈1.3 → web',
    'h3 · wget 下载 → h4:8000/download',
    '  小ACK + 海量tx数据 → 收发比≈0.19 → download',
    '',
    'AI 输入 10 维特征 (全部来自 PortStatsReply 差值):',
    '  rx/tx 包数 · rx/tx 字节 · 包速率 · 字节速率',
    '  平均包大小 · tx包速率 · 采集间隔 · 收发比',
    '核心洞察: OVS 只看到包数/字节数/收发比',
    '  AI 不关心你用 curl 还是 socket，只看底层指标',
], lsz=10)

# 右列：AI 模型 + 结果
card(s2, 9.1, 1.6, 3.8, 2.5, 'AI 模型 · Random Forest', [
    '结构: 200棵树, max_depth=15, 投票决定分类',
    '训练: scikit-learn, CPU 秒级完成',
    '推理: pickle加载, <1ms, 不需要GPU',
    '',
    '与老师案例 (ResNet-1D) 对比:',
    '  RF(我们)       |  ResNet(案例)',
    '  10维流统计      |  原始数据包',
    '  <1ms, CPU      |  需GPU加速',
    '  Flask微服务     |  Jupyter离线',
    '  在线实时系统     |  离线高精度',
    '',
    '侧重: 在线+实时+系统集成，非纯模型精度',
], lsz=10)

# 底部结果条
card(s2, 0.5, 4.3, 12.3, 1.0, '实验结果', [
    '连通性: 0% 丢包 (4主机互通)    |    AI 分类: 3/3 正确 (视频100%·网页95%·下载98%)    |    AI API: 累计 65 次 HTTP 请求',
    '对比实验: 纯L2模式下 0/3识别·0次API·无QoS  →  有AI后 3/3识别·65次API·视频VIP/网页MED/下载LOW',
], tc=GREEN, tsz=11, lsz=10)

# 底部技术细节
card(s2, 0.5, 5.5, 6.0, 1.7, 'AI 为什么在控制层 (不在数据层)', [
    '数据层 (OVS): C语言转发引擎, 不能运行Python/sklearn',
    '  · 只做两件事: 匹配流表转发 · 提供端口计数器',
    '控制层 (Ryu): 全局视野 + 编程能力, 适合跑AI',
    '  · AI从嵌入式升级为独立Flask进程: 松耦合 · 可替换',
    '  · curl就能验证: POST /predict 返回JSON分类结果',
    '  · requests_served 计数器追踪每次真实API调用',
], lsz=10)

card(s2, 6.8, 5.5, 6.0, 1.7, '核心代码', [
    '# Ryu 控制器每2秒采集端口统计',
    'features = [rx_pkts, rx_bytes, tx_pkts, tx_bytes,',
    '           pkt_rate, byte_rate, avg_pkt_size,',
    '           tx_pkt_rate, dt*1000, rx_pkts/max(tx_pkts,1)]',
    '',
    '# 通过 HTTP 调用 AI 推理服务',
    'r = requests.post("http://127.0.0.1:5001/predict",',
    '                  json={"features": features})',
    'cls = r.json()["class"]  # → video/web/download',
], lsz=10)


# ═══════════════════════════════════════════════════════
# Slide 3: 作品目标问题
# ═══════════════════════════════════════════════════════
s3 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s3)
title_bar(s3, '作品目标问题')

card(s3, 0.5, 1.2, 12.3, 1.8,
    '问题 1：传统网络不知道自己上面在跑什么类型的流量', [
    '传统交换机（纯 L2 转发）能根据 MAC 地址正确转发数据包，但对每个端口传输什么类型的流量一无所知。',
    '视频流（UDP 大包高频）和文件下载（TCP 大流量）在交换机眼里完全一样——都是「包」。',
    '没有类型感知能力 → 网络只能做无差别转发 → 无法为关键业务提供服务质量保障。',
    '',
    '对比: 纯 L2 模式下 OVS 流表只有 table-miss + MAC 转发规则 · 0 端口被识别 · 0 次 AI API 调用 · 无 ANY QoS 策略',
], tc=RED, lsz=11)

card(s3, 0.5, 3.2, 12.3, 1.8,
    '问题 2：AI 应该放在 SDN 架构的哪一层？如何设计松耦合的集成方案？', [
    '数据层的 OVS 交换机是 C 语言转发引擎，不能运行 Python、不能加载 scikit-learn 模型。AI 只能放在控制层。',
    '但如果把 AI 直接嵌入 Ryu（进程内 pickle.load + .predict()），会形成紧耦合——换模型就需要改控制器代码。',
    '而且嵌入式的 AI 无法被独立验证——老师没法用标准工具（curl）直接测试它是不是真实在运行。',
    '',
    '需要回答: AI 服务如何独立部署？Ryu 如何通过网络调用？如何用 curl 就能验证 AI 的每次推理都是真实发生的？',
], tc=RED, lsz=11)

card(s3, 0.5, 5.2, 12.3, 1.8,
    '问题 3：如何证明 AI 真的有效？——需要一个严格的对照实验', [
    '主观说「AI 能分类」没有说服力。需要一个对照实验：同一 Mininet 拓扑、完全相同的三种流量，',
    '唯一的变量是控制器有没有加载 AI（其他完全不变）。Phase 1 跑纯 L2 控制器（ryu_l2_only.py，仅 55 行代码），',
    'Phase 2 跑 AI 控制器（ryu_qos_ai.py + ai_server.py）。对比指标：流量识别端口数、AI API 调用次数、QoS 策略是否生效。',
    '',
    '同时 AI API 必须可独立调用——用 curl POST 发特征向量，直接拿到 JSON 分类结果和 request_id 计数器。',
    '预期: Phase 1 全部为 0 · Phase 2 3/3 识别正确——量化的对比结论比任何主观描述都有说服力。',
], tc=RED, lsz=11)


# ═══════════════════════════════════════════════════════
# Slide 4: 解决思路
# ═══════════════════════════════════════════════════════
s4 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s4)
title_bar(s4, '解决思路')

# 上排两列
card(s4, 0.5, 1.1, 6.1, 2.8, '方案 1：混合流量方案 + 收发比特征', [
    '三种工具 → 三种协议特征 → 收发比 (rx/tx) 区分:',
    '',
    'h1 Python socket UDP → 收发比≈∞ (无回包)',
    '  大包(1200B)·高频(~200pps)·UDP单向 → 判定为 video',
    '',
    'h2 curl 间歇 HTTP → 收发比≈1 (双向均衡)',
    '  中包(550B)·间歇(~30pps)·TCP双向 → 判定为 web',
    '',
    'h3 wget 持续下载 → 收发比≈0.19 (tx>>rx)',
    '  客户端ACK(64B)+服务器数据(1400B)·TCP单向 → download',
    '',
    'AI 输入: 10 维特征（全部来自 OVS PortStatsReply 差值）',
    '  包速率、字节速率、平均包大小、收发比、采集间隔…',
    '核心洞察: 不管上层用 curl/wget/socket，',
    'AI 只看交换机端口统计的底层指标，不关心工具类型',
], lsz=10)

card(s4, 6.9, 1.1, 6.0, 2.8, '方案 2：AI 在控制层 + 独立微服务架构', [
    'AI 放在控制层（不是数据层），这是 SDN 的正确设计:',
    '',
    '数据层 (OVS) — 不能做 AI:',
    '  C 语言转发引擎，不能运行 Python/scikit-learn',
    '  职责: 匹配流表转发 + 提供端口计数器',
    '',
    '控制层 (Ryu) — AI 推理在这里:',
    '  全局视野·Python 编程·毫秒级推理·2s 采集周期够用',
    '',
    'AI 拆为独立 Flask 微服务 (端口 5001):',
    '  Ryu → HTTP POST → Flask → predict() → JSON 结果',
    '',
    '好处: 换模型只改 Flask·控制器代码不动',
    '  curl 直接测 AI·证明不是嵌在控制器里的假模块',
    '  符合工业界 ONOS + 独立 ML 集群的架构模式',
], lsz=10)

# 下排两列
card(s4, 0.5, 4.1, 6.1, 2.8, '方案 3：对照组实验 — 量化证明 AI 的价值', [
    '同一拓扑 + 同一流量 · 唯一变量: 控制器有没有 AI',
    '',
    'Phase 1 — 纯 L2 (ryu_l2_only.py, 仅 55 行代码):',
    '  MAC学习+L2转发 · 无统计采集 · 不调AI · 无QoS',
    '  OVS流表: 仅 table-miss + MAC转发规则',
    '  结果:  0/3 端口识别 · 0 次 API · 无 QoS',
    '',
    'Phase 2 — AI (ryu_qos_ai.py + ai_server.py):',
    '  每2s采集OVS端口统计·10维特征·HTTP→AI→分类',
    '  结果:  3/3 端口识别 · 60+ 次 API · VIP/MED/LOW',
    '',
    '对比表:',
    '┌────────────┬─────────┬──────────┐',
    '│ 指标       │ 纯 L2   │ AI       │',
    '├────────────┼─────────┼──────────┤',
    '│ 流量识别   │ 0/3     │ 3/3      │',
    '│ AI API调用 │ 0       │ 65       │',
    '│ QoS        │ 无      │ 3级区分   │',
    '└────────────┴─────────┴──────────┘',
], lsz=10)

card(s4, 6.9, 4.1, 6.0, 2.8, '方案 4：AI 可独立验证 — 不是黑盒子', [
    'AI 推理服务是独立的 Flask 进程 (端口 5001)，',
    '任何人都可以用 curl 直接测试它:',
    '',
    'curl -X POST http://127.0.0.1:5001/predict \\',
    '  -H "Content-Type: application/json" \\',
    '  -d \'{"features":[400,480000,3,192,180,192000,1200,1.5,2000,100]}\'',
    '',
    '→ {"class":"video","confidence":1.0,',
    '   "probabilities":{"video":1.0,"web":0.0,"download":0.0},',
    '   "request_id":85}',
    '',
    '关键设计:',
    '  requests_served 计数器每调一次 +1',
    '  curl 就能验证 → AI 不是假模块，每次推理都是真实的',
    '  Flask 独立进程，松耦合 → 换模型不改控制器一行代码',
    '  这就是工业界「控制器 + 独立 ML 推理」的部署模式',
], lsz=10)

# 底部
txt(s4, 0.5, 7.05, 12.3, 0.3,
    '收发比区分三类流量  |  AI 从嵌入式升级为独立微服务  |  对照实验量化 AI 价值  |  curl 可验证，不是黑盒子',
    sz=11, c=BLUE, bold=True, align=PP_ALIGN.CENTER)


# ── 保存 ──
prs.save(OUTPUT)
print(f'PPT 已保存: {OUTPUT}')
print(f'共 {len(prs.slides)} 页')
